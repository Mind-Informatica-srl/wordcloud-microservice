import io
import copy
from pptx import Presentation

def copia_slide(presentation, slide):
    """ Duplica una slide copiando direttamente gli elementi XML, mantenendo contenuto, dimensioni e posizione. """
    new_slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    for shape in slide.shapes:
        new_shape = None
        if shape.shape_type == 13:  # Immagine
            image_stream = io.BytesIO(shape.image.blob)
            new_image = new_slide.shapes.add_picture(image_stream, shape.left, shape.top, shape.width, shape.height)

            # Copia il testo alternativo (alt text) dall'XML
            image_element = shape._element
            if image_element is not None:
                cNvPr = image_element.find('.//p:cNvPr', namespaces={'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
                if cNvPr is not None:
                    alt_text = cNvPr.get('descr')
                    if alt_text:
                        new_image._element.find('.//p:cNvPr', namespaces={'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}).set('descr', alt_text)
        else:
            # Copia diretta dell'elemento XML della forma
            el = shape.element
            new_shape = copy.deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(new_shape, 'p:extLst')

    return new_slide

# def filtra_per(ppt, replacements):
#     slides_to_keep = []

#     for count, slide in enumerate(ppt.slides):
#         keep_slide = True  # Supponiamo di mantenerla

#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 for para in shape.text_frame.paragraphs:
#                     if para.text.startswith("{{if:"):
#                         conditions = [cond.split("}}")[0].strip() for cond in para.text.split("{{if:") if "}}" in cond]
#                         for condition in conditions:
#                             parts = condition.split(":")
#                             if len(parts) != 2:
#                                 continue

#                             ph = "{{" + parts[0] + "}}"
#                             v = parts[1]

#                             if v == "*":
#                                 if ph not in replacements or not replacements[ph] or replacements[ph] == "":
#                                     keep_slide = False
#                                     break
#                                 elif replacements[ph] and replacements[ph] != "":
#                                     para.text = para.text.replace("{{if:" + condition + "}}", "")
#                             elif v == "":
#                                 if ph in replacements.keys() and replacements[ph] == "":
#                                     keep_slide = False
#                                     break
#                                 elif ph not in replacements.keys():
#                                     para.text = para.text.replace("{{if:" + condition + "}}", "")
#                             else:
#                                 if replacements[ph] != v:
#                                     keep_slide = False
#                                     break
#                                 elif replacements[ph] == v:
#                                     para.text = para.text.replace("{{if:" + condition + "}}", "")

#         if keep_slide:
#             slides_to_keep.append(slide)

#     # Creiamo una nuova presentazione
#     new_ppt = Presentation()

#     for slide in slides_to_keep:
#         copia_slide(new_ppt, slide)

#     return new_ppt  # Restituiamo la nuova presentazione

def filtra_per(ppt, replacements):
    """
        controllo tutte le slide, se la slide possiede la stringa {{if:val:*}}
        se non esiste vado avanti
        se esiste controllo se la stringa è uguale a quella passata
        se è uguale la tengo
        se non è uguale la elimino
    """
    slides_to_remove = []

    for count, slide in enumerate(ppt.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.startswith("{{if:"):
                        conditions = [cond.split("}}")[0].strip() for cond in para.text.split("{{if:") if "}}" in cond]
                        for condition in conditions:
                            t = condition
                            parts = t.split(":")
                            if len(parts) == 2:
                                ph = "{{" + parts[0] + "}}"
                                v = parts[1]
                            else:
                                continue

                            if v == "*":
                                if ph not in replacements.keys() or replacements[ph] is None or replacements[ph] == '':
                                    slides_to_remove.append(count)
                                    break
                                elif replacements[ph] is not None and replacements[ph] != '':
                                    para.text = para.text.replace("{{if:" + condition + "}}", "")
                            elif v == "":
                                if ph in replacements.keys() and replacements[ph] == '':
                                    slides_to_remove.append(count)
                                    break
                                elif ph not in replacements.keys() or (ph in replacements.keys() and replacements[ph] != ''):
                                    para.text = para.text.replace("{{if:" + condition + "}}", "")
                            else:
                                if replacements[ph] != v:
                                    slides_to_remove.append(count)
                                    break
                                elif replacements[ph] == v:
                                    para.text = para.text.replace("{{if:" + condition + "}}", "")

    # Elimina le slide raccolte
    for i in range(len(slides_to_remove)-1, -1, -1):
        index = slides_to_remove[i]
        rId = ppt.slides._sldIdLst[index].rId
        ppt.part.drop_rel(rId)
        del ppt.slides._sldIdLst[index]

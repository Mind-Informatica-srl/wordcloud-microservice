import copy
from pptx import Presentation
import docx
from docx import Document
import os
import shutil
from python_pptx_text_replacer import TextReplacer
from PIL import Image

from constants import UPLOAD_FOLDER
from office.duplicate_and_replace_slide import duplicate_and_replace_slide, extract_placeholder_info
from office.filtra_per import filtra_per
from office.save_image import save_image
import re
import io
import subprocess
import requests
from docx.shared import Inches
from io import BytesIO
import base64

import concurrent.futures

from pdf2image import convert_from_bytes


from python_docx_replace import docx_blocks, docx_replace
from python_docx_replace.paragraph import Paragraph


def replace_image_in_pptx(ppt, image_replacements):
    # Sostituire le immagini
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # 13 è il tipo per l'immagine
                image_element = shape._element

                # Cerca il testo alternativo nell'elemento XML (se presente)
                alt_text = None
                if image_element is not None:
                    # Cerca il testo alternativo all'interno dell'elemento XML
                    cNvPr = image_element.find('.//p:cNvPr', namespaces={'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
                    if cNvPr is not None:
                        alt_text = cNvPr.get('descr')



                for placeholder, image_path in image_replacements.items():
                    if placeholder in shape.image.filename:
                        # Posizione e dimensioni dell'immagine
                        left, top, or_width, or_height = shape.left, shape.top, shape.width, shape.height
                        # Aggiungere la nuova immagine
                        imaga_saved = save_image(placeholder, image_path, "storage/immagini", width=or_width, height=or_height)
                        dimfissa = or_width if or_width > or_height else or_height

                        if placeholder.startswith("{{C5}}") or placeholder.startswith("{{C6}}") or placeholder.startswith("{{C7}}"):
                                            height = or_height
                                            width = or_width
                        else:
                            with Image.open(imaga_saved[placeholder]) as img:
                                new_width, new_height = img.size

                            # Calcolare le nuove dimensioni mantenendo il rapporto
                            if dimfissa == or_width:
                                width = or_width
                                height = int(new_height * (or_width / new_width))
                            else:
                                height = or_height
                                width = int(new_width * (or_height / new_height))

                            if width > or_width or height > or_height:
                                # Ridimensiona l'immagine per adattarla
                                if width > or_width:
                                    width = or_width
                                    height = int(new_height * (or_width / new_width))
                                elif height > or_height:
                                    height = or_height
                                    width = int(new_width * (or_height / new_height))

                        slide.shapes.add_picture(imaga_saved[placeholder], left, top, width, height)
                    elif alt_text is not None and placeholder in alt_text:
                        # Posizione e dimensioni dell'immagine
                        left, top, or_width, or_height = shape.left, shape.top, shape.width, shape.height
                        # Aggiungere la nuova immagine
                        imaga_saved = save_image(placeholder, image_path, "storage/immagini", width=or_width, height=or_height)
                        dimfissa = or_width if or_width > or_height else or_height

                        if placeholder.startswith("{{C5}}") or placeholder.startswith("{{C6}}") or placeholder.startswith("{{C7}}"):
                                            height = or_height
                                            width = or_width
                        else:
                            with Image.open(imaga_saved[placeholder]) as img:
                                new_width, new_height = img.size

                            # Calcolare le nuove dimensioni mantenendo il rapporto
                            if dimfissa == or_width:
                                width = or_width
                                height = int(new_height * (or_width / new_width))
                            else:
                                height = or_height
                                width = int(new_width * (or_height / new_height))

                            if width > or_width or height > or_height:
                                # Ridimensiona l'immagine per adattarla
                                if width > or_width:
                                    width = or_width
                                    height = int(new_height * (or_width / new_width))
                                elif height > or_height:
                                    height = or_height
                                    width = int(new_width * (or_height / new_height))

                        slide.shapes.add_picture(imaga_saved[placeholder], left, top, width, height)
                        # Rimuove l'immagine originale
                        sp = shape
                        slide.shapes._spTree.remove(sp._element)

def salva_byte_pptx(ppt):
    # Salvare il file PPTX modificato
    storage_path = os.path.join("storage", "modificati")
    os.makedirs(storage_path, exist_ok=True)
    modified_pptx_path = os.path.join("storage", "modificati", "modified_pptx.pptx")
    ppt.save(modified_pptx_path)
    print(f"File PPTX modificato e salvato come '{modified_pptx_path}'")
    # Restituire i byte del file PPTX modificato
    with open(modified_pptx_path, "rb") as f:
        pptx_bytes = f.read()

    return pptx_bytes


def delete_paragraph(paragraph):
    p = paragraph._element
    p.getparent().remove(p)
    p._p = p._element = None

def insert_pdf_images_into_docx(doc, pdf_bytes, insert_index):
    images = convert_from_bytes(pdf_bytes)
    for img in images:
        # Converti immagine in PNG e salvala in un buffer
        buffer = BytesIO()
        img.save(buffer, format='PNG')
        buffer.seek(0)

        # Inserisci paragrafo e immagine nel documento
        doc.paragraphs[insert_index].insert_paragraph_before()
        doc.paragraphs[insert_index - 1].add_run().add_picture(buffer, width=Inches(6))

    return insert_index + len(images)

def append_to_doc(doc,p):
    doc.add_paragraph("",p.style)       # add an empty paragraph in the matching style
    for r in p.runs:
        nr = doc.paragraphs[-1].add_run(r.text)
        nr.style = r.style
        nr.bold = r.bold
        nr.italic = r.italic
        nr.underline = r.underline


def duplica_blocchi_paragrafi(doc, replacements_for_each):
    i = 0
    while i < len(doc.paragraphs):
        para = doc.paragraphs[i]
        # if para.text.startswith("{{for_fg:") or para.text.startswith("{{for_go:"):
        if "{{for_fg:" in para.text or "{{for_go:" in para.text:
            if "{{for_fg:" in para.text:
                forcontroller = "{{fine_for_fg}}"
                forplaceholder = "{{for_fg:n}}"
            else:
                forcontroller = "{{fine_for_go}}"
                forplaceholder = "{{for_go:n}}"
            print("Trovato inizio blocco di duplicazione", i)
        match = re.match(r"{{for_(fg|go):(\d+)}}", para.text.strip())
        if match:
            n = int(match.group(2))
            start_idx = i
            # Trova la fine del blocco
            end_idx = None
            for j in range(i+1, len(doc.paragraphs)):
                # if doc.paragraphs[j].text.endswith(forcontroller):
                if forcontroller in doc.paragraphs[j].text:
                    end_idx = j
                    break
            if end_idx is None:
                i += 1
                continue  # Nessuna fine trovata, ignora

            # Copia i paragrafi del blocco (esclude marker)
            blocco = [doc.paragraphs[k] for k in range(start_idx+1, end_idx)]
            # Duplicazione
            replacements = replacements_for_each[forplaceholder]
            nuovi_paragrafi = []
            for dup_idx in range(n):
                rep = replacements[dup_idx] if dup_idx < len(replacements) else replacements[0]
                # Ottieni i replacements per questa duplicazione (se presenti)
                for p in blocco:
                    new_par = doc.add_paragraph("", p.style)
                    for r in p.runs:
                        nr = new_par.add_run(r.text)
                        nr.style = r.style
                        nr.bold = r.bold
                        nr.italic = r.italic
                        nr.underline = r.underline
                    paragraph_obj = Paragraph(new_par)
                    for key, value in rep["testuali"].items():
                        paragraph_obj.replace_key(key, str(value))
                    # Aggiorna eventuali immagini nel paragrafo
                    # Copia solo le immagini associate a questo paragrafo
                    for shape in doc.inline_shapes:
                        # Verifica se l'immagine è nel paragrafo corrente
                        parent = shape._inline.getparent()
                        while parent is not None and parent != p._element:
                            parent = parent.getparent()
                        if parent == p._element:
                            alt_text = shape._inline.docPr.get('descr')
                            if alt_text:
                                new_alt = re.sub(
                                    r"(\{\{[a-zA-Z0-9_]+:)\d+(\}\})",
                                    lambda m: f"{m.group(1)}{dup_idx + 2}{m.group(2)}",
                                    alt_text
                                )
                            else:
                                new_alt = None
                            
                            info = extract_placeholder_info(new_alt)
                            if info is not None:
                                for placeholder, image_path in rep["immagini"].items():
                                    if placeholder in info:
                                        imaga_saved = save_image(placeholder, image_path, "storage/immagini", width=0, height=0)
                                        new_alt = ""
                                        
                                width, height = shape.width, shape.height
                                # blip = shape._inline.graphic.graphicData.pic.blipFill.blip
                                with open(imaga_saved[info[1]], "rb") as img_file:
                                     new_image_data = img_file.read()

                                run_img = new_par.add_run()
                                run_img.add_picture(io.BytesIO(new_image_data), width=width, height=height)
                                # image_part = doc.part.related_parts[blip.embed]
                                # image_part._blob = new_image_data
                                last_shape = doc.inline_shapes[-1]
                                if new_alt:
                                    last_shape._inline.docPr.set('descr', new_alt)
                            else:
                                width, height = shape.width, shape.height
                                blip = shape._inline.graphic.graphicData.pic.blipFill.blip
                                image_part = doc.part.related_parts[blip.embed]
                                original_image_data = image_part.blob

                                run_img = new_par.add_run()
                                run_img.add_picture(io.BytesIO(original_image_data), width=width, height=height)
                                # Copia anche l'alt text originale, se presente
                                if alt_text:
                                    last_shape = doc.inline_shapes[-1]
                                    last_shape._inline.docPr.set('descr', alt_text)
                    nuovi_paragrafi.append(new_par)
            # Rimuovi il blocco originale (compresi i marker)
            for idx in range(end_idx, start_idx - 1, -1):
                delete_paragraph(doc.paragraphs[idx])

            # Inserisci i nuovi paragrafi al posto del blocco originale
            body = doc._body._element
            insert_pos = start_idx
            for new_par in nuovi_paragrafi:
                body.insert(insert_pos, new_par._element)
                insert_pos += 1
            # Aggiorna l'indice per saltare i duplicati appena inseriti
            i = start_idx + len(nuovi_paragrafi)

        else:
            i += 1

def replace_text_in_docx(docx_path, replacements, image_replacements, replacements_for_each):
    # print_runs_in_docx(docx_path)
    # Caricare il file DOCX
    doc = docx.Document(docx_path)
    docx_replace(doc, **replacements)
    duplica_blocchi_paragrafi(doc, replacements_for_each)
    # Sostituire le immagini tramite testo alternativo
    doc, mantieni_idx, rimuovi_idx = valuta_if_docx(doc, replacements, image_replacements)
    # elabora_blocchi_paragrafi(doc, replacements_for_each, replacements, image_replacements)
    docx_blocks(doc, mantieni_idx, rimuovi_idx, da_mantenere=True, da_rimuovere=False)  
    replace_image_in_docx(doc, image_replacements)
    # Per ogni replacement_for_each, recuperare i valori del campo Html
    if replacements_for_each is not None and len(replacements_for_each) > 0:
        if "{{for_go:n}}" in replacements_for_each:
            # se esiste, cerco all'interno del file docx il blocco che inizia con {{liste_controllo_go:}} e termina con {{fine_liste_controllo_go:}}
            # recupero i byte eseguendo la richiesta al server con la url contenuta in item["html"]
            # per ogni elemento contenuto in {{for_go:n}} sostituisco il blocco {{liste_controllo_go:}} con i byte del file docx recuperato
              # e sostituisco {{fine_liste_controllo_go:}} con un tag <fine_liste_controllo_go>

             for idx, para in enumerate(doc.paragraphs):
                 if "{{liste_controllo_go}}" in para.text:
                    for item in replacements_for_each["{{for_go:n}}"]:
                        if "html" in item and item["html"] is not None:
                            url = os.getenv("BASE_URL") + item["html"]["{{liste_controllo_go}}"]
                            response = requests.get(url)
                            if response.status_code == 200:
                                pdf_bytes = response.content
                                insert_pdf_images_into_docx(doc, pdf_bytes, idx)
                            else:
                                print(f"Errore durante il download da {url}")
                    para.text = para.text.replace("{{liste_controllo_go}}", "")
                    break  # Rimuovi se vuoi inserire lo stesso blocco in più punti
                    # i = 0
                    # for para in doc.paragraphs:
                    #     if para.text.startswith("{{liste_controllo_go}}"):
                    #         # Trova la fine del blocco
                    #         end_idx = None
                    #         for j in range(i + 1, len(doc.paragraphs)):
                    #             if doc.paragraphs[j].text.strip().startswith("{{fine_liste_controllo_go}}"):
                    #                 end_idx = j
                    #                 break
                    #         if end_idx is not None:
                    #             # Recupera i byte del file HTML
                    #             url =  os.getenv("BASE_URL") + item["html"]["{{liste_controllo_go}}"]
                    #             response = requests.get(url)
                    #             if response.status_code == 200:
                    #                 html_bytes = response.content
                    #             else:
                    #                 print(f"Errore durante il download dell'immagine da '{url}'")
                    #             # Crea un nuovo documento DOCX con i byte HTML
                    #             new_doc = docx_to_pdf_merge(doc, html_bytes)
                    #             # Aggiungi il nuovo documento al documento principale
                    #             for new_para in new_doc.paragraphs:
                    #                 doc.add_paragraph(new_para.text, style=new_para.style)
                    #             # Rimuovi il blocco originale
                    #             for idx in range(end_idx, doc.paragraphs.index(para), -1):
                    #                 delete_paragraph(doc.paragraphs[idx])
                    #             delete_paragraph(para)
                    #     else:
                    #         i += 1


    # Salvare il file DOCX modificato
    storage_path = os.path.join("storage", "modificati")
    os.makedirs(storage_path, exist_ok=True)
    modified_docx_path = os.path.join("storage", "modificati", "modified_docx.docx")
    doc.save(modified_docx_path)
    print(f"File DOCX modificato e salvato come '{modified_docx_path}'")

    # Restituire i byte del file PPTX modificato
    with open(modified_docx_path, "rb") as f:
        docx_bytes = f.read()

    return docx_bytes

def print_runs_in_docx(docx_path):
    doc = Document(docx_path)
    for i, paragraph in enumerate(doc.paragraphs):
        print(f"Paragrafo {i} XML completo:\n{paragraph._element.xml}\n")

def aggiorna_toc_con_libreoffice(input_path, output_path=None):
    """
    Aggiorna automaticamente il sommario di un file .docx usando LibreOffice in modalità headless.
    Se output_path non è specificato, il file aggiornato verrà salvato nella stessa cartella di input.
    """
    if output_path is None:
        output_dir = os.path.dirname(input_path)
    else:
        output_dir = os.path.dirname(output_path)
    # Comando per convertire e aggiornare i campi (incluso il TOC)
    cmd = [
        "soffice",
        "--headless",
        "--convert-to", "docx",
        "--outdir", output_dir,
        input_path
    ]
    subprocess.run(cmd, check=True)
    # Il file convertito avrà lo stesso nome base
    base_name = os.path.basename(input_path)
    converted_path = os.path.join(output_dir, base_name.replace(".docx", ".docx"))
    # Se serve, rinomina il file di output
    if output_path and converted_path != output_path:
        os.rename(converted_path, output_path)
    return output_path or converted_path

def pixel_to_emu(pixels, dpi=96):
    # 1 pollice = 914400 EMU, 1 pollice = dpi pixel
    return int(pixels * 914400 / dpi)

def bytes_to_base64(bytes_data):
    return base64.b64encode(bytes_data).decode('utf-8')

def replace_image_in_docx(doc, image_replacements):
    def process_shape(shape):
        import docx
        from PIL import Image
        alt_text = shape._inline.docPr.get('descr')
        if alt_text is not None:
            for placeholder, image_path in image_replacements.items():
                if placeholder in alt_text:
                    or_width, or_height = shape.width, shape.height
                    if placeholder.startswith("{{B") or placeholder.startswith("{{DA31_"):
                        imaga_saved = save_image(placeholder, image_path, "storage/immagini", width=or_width, height=or_height)
                    else:
                        imaga_saved = save_image(placeholder, image_path, "storage/immagini", width=0, height=0)
                    dimfissa = or_width if or_width > or_height else or_height

                    if placeholder.startswith("{{C5}}") or placeholder.startswith("{{C6}}") or placeholder.startswith("{{C7}}"):
                        height = or_height
                        width = or_width
                    else:
                        with Image.open(imaga_saved[placeholder]) as img:
                            new_width, new_height = img.size
                        new_width = pixel_to_emu(new_width)
                        new_height = pixel_to_emu(new_height)
                        if dimfissa == or_width:
                            width = or_width
                            height = int(new_height * (or_width / new_width))
                        else:
                            height = or_height
                            width = int(new_width * (or_height / new_height))
                        if width > or_width or height > or_height:
                            if width > or_width:
                                width = or_width
                                height = int(new_height * (or_width / new_width))
                            elif height > or_height:
                                height = or_height
                                width = int(new_width * (or_height / new_height))
                    blip = shape._inline.graphic.graphicData.pic.blipFill.blip
                    with open(imaga_saved[placeholder], "rb") as new_image_file:
                        new_image_data = new_image_file.read()
                    image_part = doc.part.related_parts[blip.embed]
                    image_part._blob = new_image_data

    # Parallelizza sulle immagini
    with concurrent.futures.ThreadPoolExecutor() as executor:
        list(executor.map(process_shape, doc.inline_shapes))

def replace_image_in_docx1(doc, image_replacements):
    #######################
    for shape in doc.inline_shapes:
        if shape.type == docx.enum.shape.WD_INLINE_SHAPE.PICTURE or shape.type == 3:
            alt_text = shape._inline.docPr.get('descr')
            if alt_text is not None:
                for placeholder, image_path in image_replacements.items():
                    if placeholder in alt_text:
                        or_width, or_height = shape.width, shape.height
                        if placeholder.startswith("{{B") or placeholder.startswith("{{DA31_") or placeholder.startswith("{{C5") or placeholder.startswith("{{C7"):
                            imaga_saved = save_image(placeholder, image_path, "storage/immagini", width=or_width, height=or_height)
                        else:
                            imaga_saved = save_image(placeholder, image_path, "storage/immagini", width=0, height=0)
                        # Recupera il "blip" (che contiene il riferimento all'immagine nel pacchetto)
                        # blip = shape._inline.graphic.graphicData.pic.blipFill.blip

                        dimfissa = or_width if or_width > or_height else or_height

                        if placeholder.startswith("{{C5}}") or placeholder.startswith("{{C7}}"):
                                            height = or_height
                                            width = or_width
                        else:
                            with Image.open(imaga_saved[placeholder]) as img:
                                new_width, new_height = img.size

                            new_width = pixel_to_emu(new_width)
                            new_height = pixel_to_emu(new_height)

                            # Calcolare le nuove dimensioni mantenendo il rapporto
                            if dimfissa == or_width:
                                width = or_width
                                height = int(new_height * (or_width / new_width))
                            else:
                                height = or_height
                                width = int(new_width * (or_height / new_height))

                            if width > or_width or height > or_height:
                                # Ridimensiona l'immagine per adattarla
                                if width > or_width:
                                    width = or_width
                                    height = int(new_height * (or_width / new_width))
                                elif height > or_height:
                                    height = or_height
                                    width = int(new_width * (or_height / new_height))

                        # Recupera il "blip" (che contiene il riferimento all'immagine nel pacchetto)
                        blip = shape._inline.graphic.graphicData.pic.blipFill.blip

                        with open(imaga_saved[placeholder], "rb") as new_image_file:
                            new_image_data = new_image_file.read()
                        # bb = bytes_to_base64(new_image_data)
                        # print(f"Replacing image for placeholder '{placeholder}' with base64 data of length {len(bb)}")
                        # print(blip.embed)
                        image_part = doc.part.related_parts[blip.embed]
                        image_part._blob = new_image_data


                        # # Rimuove immagine esistente
                        # drawing_element = shape._inline
                        # drawing_parent = drawing_element.getparent()
                        # drawing_parent.remove(drawing_element)
                        # # Trova il paragrafo corrispondente e aggiunge la nuova immagine
                        # paragraph = None
                        # for para in doc.paragraphs:
                        #     if shape._inline in [r._element for r in para.runs]:
                        #         paragraph = para
                        #         break
                            
                        # if paragraph:
                        #     run = paragraph.add_run()
                        #     run.add_picture(image_path, width=width, height=height)
    ##########################
    # for paragraph in doc.paragraphs:
    #     for run in paragraph.runs:
    #         drawing_elements = run._element.xpath('.//pic:cNvPr[@descr]')
    #         if drawing_elements:
    #             descr_elem = drawing_elements[0]
    #             alt_text = descr_elem.get("descr")

    #             for placeholder, image_url in image_replacements.items():
    #                 if placeholder in alt_text:
    #                     # Dimensioni originali (approssimate)
    #                     shape = next((s for s in doc.inline_shapes if s._inline.docPr.get('descr') == alt_text), None)
    #                     if shape:
    #                         or_width = shape.width
    #                         or_height = shape.height
    #                     else:
    #                         or_width = Inches(2)
    #                         or_height = Inches(2)

    #                     # Scarica immagine e salvala
    #                     saved = save_image(placeholder, image_url, "storage/immagini", width=or_width, height=or_height)
    #                     image_path = saved[placeholder]

    #                     # Calcola proporzioni
    #                     if placeholder.startswith("{{C5}}") or placeholder.startswith("{{C6}}") or placeholder.startswith("{{C7}}"):
    #                         width = or_width
    #                         height = or_height
    #                     else:
    #                         with Image.open(image_path) as img:
    #                             new_width, new_height = img.size
    #                         dimfissa = max(or_width, or_height)

    #                         if dimfissa == or_width:
    #                             width = or_width
    #                             height = int(new_height * (or_width / new_width))
    #                         else:
    #                             height = or_height
    #                             width = int(new_width * (or_height / new_height))

    #                         # Limita alle dimensioni originali
    #                         if width > or_width:
    #                             width = or_width
    #                             height = int(new_height * (or_width / new_width))
    #                         if height > or_height:
    #                             height = or_height
    #                             width = int(new_width * (or_height / new_height))

    #                     # Rimuove il run esistente (contenente l'immagine)
    #                     paragraph._element.remove(run._element)

    #                     # Inserisce nuova immagine nella stessa posizione
    #                     new_run = paragraph.add_run()
    #                     new_run.add_picture(image_path, width=width, height=height)


def elimina_cartella(path):
    if os.path.exists(path):
        shutil.rmtree(path)

def valuta_if_docx(doc, replacements, replacements_image):
    stack = []
    blocchi = []

    # Trova tutti i blocchi if annidati
    for idx, para in enumerate(doc.paragraphs):
        text = ''.join(run.text for run in para.runs)

        for match_if in re.finditer(r"\{\{if:([^\}:]+):([^\}]*)\}\}", text):
            placeholder, condizione = match_if.groups()
            stack.append((placeholder, condizione, idx))

        for match_fine in re.finditer(r"\{\{fine_if:([^\}:]+):([^\}]*)\}\}", text):
            placeholder, condizione = match_fine.groups()
            if stack and stack[-1][:2] == (placeholder, condizione):
                start_ph, cond, start_idx = stack.pop()
                blocchi.append((start_idx, idx, start_ph, cond))
            else:
                print(f"[WARNING] Errore: fine_if senza if corrispondente a riga {idx} per {placeholder} con condizione {condizione}")

    if stack:
        for ph, cond, idx in stack:
            print(f"[WARNING] Errore: if aperto a riga {idx} senza fine_if corrispondente per {ph} con condizione {cond}")

    mantieni_idx = 1
    rimuovi_idx = 1
    # Applica modifiche (dall'interno all'esterno)
    for start_idx, end_idx, placeholder, condizione in reversed(blocchi):
        ph = f"{{{{{placeholder}}}}}"
        soddisfatta = False

        if condizione == "*":
            soddisfatta = (ph in replacements and replacements[ph] not in [None, ""]) or \
                          (ph in replacements_image and replacements_image[ph] not in [None, ""])
        elif condizione == "":
            soddisfatta = (ph in replacements and replacements[ph] == "") or \
                          (ph in replacements_image and replacements_image[ph] == "")
        else:
            soddisfatta = (ph in replacements and replacements[ph] == condizione) or \
                          (ph in replacements_image and replacements_image[ph] == condizione)

        if soddisfatta:
            tag_open = f"<da_mantenere_{mantieni_idx}>"
            tag_close = f"</da_mantenere_{mantieni_idx}>"
            mantieni_idx += 1
        else:
            tag_open = f"<da_rimuovere_{rimuovi_idx}>"
            tag_close = f"</da_rimuovere_{rimuovi_idx}>"
            rimuovi_idx += 1

        # Sostituisci solo il testo nel paragrafo (mantenendo tutto il resto)
        _sostituisci_placeholder_in_paragrafo(doc.paragraphs[start_idx],
                                              r"\{\{if:" + re.escape(placeholder) + ":" + re.escape(condizione) + r"\}\}",
                                              tag_open)
        _sostituisci_placeholder_in_paragrafo(doc.paragraphs[end_idx],
                                              r"\{\{fine_if:" + re.escape(placeholder) + ":" + re.escape(condizione) + r"\}\}",
                                              tag_close)
        
    return doc, mantieni_idx, rimuovi_idx


def _sostituisci_placeholder_in_paragrafo(paragraph, pattern_regex, replacement):
    """
    Sostituisce un segnaposto nel testo completo di un paragrafo (composto dai run),
    mantenendo gli altri contenuti e formattazioni ove possibile.
    """
    full_text = ''.join(run.text for run in paragraph.runs)
    new_text, num_subs = re.subn(pattern_regex, replacement, full_text)

    if num_subs == 0 and re.search(pattern_regex, full_text):
        print(f"[ERROR] Sostituzione fallita per pattern: {pattern_regex}")
        print(f"[ERROR] Testo paragrafo: {full_text}")

    # Cancella tutti i run e sostituisci con uno nuovo
    for run in paragraph.runs:
        run.text = ""
    if paragraph.runs:
        paragraph.runs[0].text = new_text
    else:
        paragraph.add_run(new_text)

def valuta_if_docx_in(doc, replacements, replacements_image):
    """
    Valuta se il documento contiene blocchi di paragrafi da rimuovere o mantenere.
    Se il documento contiene blocchi di paragrafi che iniziano con {{if:placeholder:condition}}
    e terminano con {{fine_if:placeholder}}, verifica se la condizione è soddisfatta.
    se la condizione non è soddisfatta, sostituisce {{if:placeholder:condition}} e {{fine_if:placeholder}} con 
    dei tag <da_rimuovere> e </da_rimuovere> per poterli rimuovere successivamente.
    Se la condizione è soddisfatta, sostituisce {{if:placeholder:condition}} con
    <da_mantenere> e </da_mantenere> per poterli mantenere successivamente.
    """
    i = 0
    while i < len(doc.paragraphs):
        para = doc.paragraphs[i]
        # if para.text.startswith("{{if:"):
        if "{{if:" in para.text:
            # Estrai placeholder e condizione
            match = re.search(r"\{\{if:([^\}:]+):([^\}]*)\}\}", para.text)
            if not match:
                i += 1
                continue
            placeholder, condizione = match.group(1), match.group(2)
            ph = "{{" + placeholder + "}}"
            # Trova la fine del blocco
            end_idx = None
            for j in range(i+1, len(doc.paragraphs)):
                # if doc.paragraphs[j].text.endswith("{{fine_if:" + placeholder + ":" + condizione + "}}"):
                if "{{fine_if:" + placeholder + ":" + condizione + "}}" in doc.paragraphs[j].text:
                    end_idx = j
                    break
            if end_idx is None:
                i += 1
                continue  # Nessuna fine trovata, ignora

            # Valuta la condizione
            soddisfatta = False
            if condizione == "*":
                soddisfatta = (ph in replacements and replacements[ph] not in [None, ""]) or (ph in replacements_image and replacements_image[ph] not in [None, ""])
            elif condizione == "":
                soddisfatta = (ph in replacements and replacements[ph] == "") or (ph in replacements_image and replacements_image[ph] == "")
            else:
                soddisfatta = (ph in replacements and replacements[ph] == condizione) or (ph in replacements_image and replacements_image[ph] == condizione)

            # Sostituisci i marker
            if soddisfatta:
                doc.paragraphs[i].text = "<da_mantenere>"
                doc.paragraphs[end_idx].text = "</da_mantenere>"
            else:
                doc.paragraphs[i].text = "<da_rimuovere>"
                doc.paragraphs[end_idx].text = "</da_rimuovere>"
            i = end_idx + 1
        else:
            i += 1

def valuta_if_docx2(doc, replacements, replacements_image):
    """
    Valuta se il documento contiene blocchi condizionali con {{if:...:...}} e {{fine_if:...:...}}.
    Sostituisce con tag <da_mantenere> o <da_rimuovere> per gestirli successivamente.
    """
    i = 0
    while i < len(doc.paragraphs):
        para = doc.paragraphs[i]
        full_text = ''.join(run.text for run in para.runs)

        if "{{if:" in full_text:
            match = re.search(r"\{\{if:([^\}:]+):([^\}]*)\}\}", full_text)
            if not match:
                i += 1
                continue
            placeholder, condizione = match.group(1), match.group(2)
            ph = "{{" + placeholder + "}}"

            # Trova l'indice del paragrafo con {{fine_if:...}}
            end_idx = None
            for j in range(i+1, len(doc.paragraphs)):
                end_text = ''.join(run.text for run in doc.paragraphs[j].runs)
                if f"{{{{fine_if:{placeholder}:{condizione}}}}}" in end_text:
                    end_idx = j
                    break

            if end_idx is None:
                i += 1
                continue  # Nessuna fine trovata

            # Valuta condizione
            soddisfatta = False
            if condizione == "*":
                soddisfatta = (ph in replacements and replacements[ph] not in [None, ""]) or (ph in replacements_image and replacements_image[ph] not in [None, ""])
            elif condizione == "":
                soddisfatta = (ph in replacements and replacements[ph] == "") or (ph in replacements_image and replacements_image[ph] == "")
            else:
                soddisfatta = (ph in replacements and replacements[ph] == condizione) or (ph in replacements_image and replacements_image[ph] == condizione)

            # Sostituzione nei paragrafi
            if soddisfatta:
                doc.paragraphs[i].text = "<da_mantenere>" + doc.paragraphs[i].text
                doc.paragraphs[end_idx].text = doc.paragraphs[i].text + "</da_mantenere>"
            else:
                doc.paragraphs[i].text = "<da_rimuovere>" + doc.paragraphs[i].text
                doc.paragraphs[end_idx].text = doc.paragraphs[i].text + "</da_rimuovere>"

            i = end_idx + 1
        else:
            i += 1


def valuta_if_blocco(blocco, replacements):
    """
    Valuta se il blocco contiene paragrafi da rimuovere o mantenere.
    Sostituisce i marker {{if:placeholder:condition}} e {{fine_if:placeholder}}
    con <da_mantenere>/<da_rimuovere> e </da_mantenere>/</da_rimuovere> in base alla condizione.
    """
    i = 0
    while i < len(blocco):
        para = blocco[i]
        if para.text.startswith("{{if:"):
            match = re.search(r"\{\{if:([^\}:]+):([^\}]*)\}\}", para.text)
            if not match:
                i += 1
                continue
            placeholder, condizione = match.group(1), match.group(2)
            ph = "{{" + placeholder + "}}"
            # Trova la fine del blocco
            end_idx = None
            for j in range(i+1, len(blocco)):
                if blocco[j].text.endswith("{{fine_if:" + placeholder + ":" + condizione + "}}"):
                    end_idx = j
                    break
            if end_idx is None:
                i += 1
                continue

            # Valuta la condizione
            soddisfatta = False
            if condizione == "*":
                soddisfatta = ph in replacements and replacements[ph] not in [None, ""]
            elif condizione == "":
                soddisfatta = ph in replacements and replacements[ph] == ""
            else:
                soddisfatta = ph in replacements and replacements[ph] == condizione

            # Sostituisci i marker
            if soddisfatta:
                blocco[i].text = "<da_mantenere>"
                blocco[end_idx].text = "</da_mantenere>"
            else:
                blocco[i].text = "<da_rimuovere>"
                blocco[end_idx].text = "</da_rimuovere>"
            i = end_idx + 1
        else:
            i += 1

    return blocco

def elabora_blocchi_paragrafi(doc, replacements_for_each, replacements, image_replacements):
    i = 0
    while i < len(doc.paragraphs):
        para = doc.paragraphs[i]
        testo = para.text.strip()

        # -------------------------------
        # GESTIONE BLOCCO {{for_fg:n}} O {{for_go:n}}
        # -------------------------------
        match_for = re.match(r"\{\{for_(fg|go):(\d+)\}\}", testo)
        if match_for:
            tipo, n = match_for.group(1), int(match_for.group(2))
            forcontroller = f"{{{{fine_for_{tipo}}}}}"
            forplaceholder = f"{{{{for_{tipo}:n}}}}"
            start_idx = i

            end_idx = next(
                (j for j in range(i+1, len(doc.paragraphs)) if forcontroller in doc.paragraphs[j].text),
                None
            )
            if end_idx is None:
                i += 1
                continue

            blocco = [doc.paragraphs[k] for k in range(start_idx+1, end_idx)]
            nuovi_paragrafi = []
            reps = replacements_for_each.get(forplaceholder, [])

            for dup_idx in range(n):
                rep = reps[dup_idx] if dup_idx < len(reps) else reps[0]

                for p in blocco:
                    new_par = doc.add_paragraph("", p.style)
                    for r in p.runs:
                        nr = new_par.add_run(r.text)
                        nr.style = r.style
                        nr.bold = r.bold
                        nr.italic = r.italic
                        nr.underline = r.underline

                    # Sostituzione testuale
                    paragraph_obj = Paragraph(new_par)
                    for key, value in rep["testuali"].items():
                        paragraph_obj.replace_key(key, str(value))

                    # Copia immagini associate
                    for shape in doc.inline_shapes:
                        parent = shape._inline.getparent()
                        while parent is not None and parent != p._element:
                            parent = parent.getparent()

                        if parent == p._element:
                            alt_text = shape._inline.docPr.get('descr')
                            width, height = shape.width, shape.height

                            # Aggiorna alt text con indice duplica
                            if alt_text:
                                new_alt = re.sub(
                                    r"(\{\{[a-zA-Z0-9_]+:)\d+(\}\})",
                                    lambda m: f"{m.group(1)}{dup_idx + 2}{m.group(2)}",
                                    alt_text
                                )
                            else:
                                new_alt = None

                            info = extract_placeholder_info(new_alt)

                            # Se l’immagine è legata a un placeholder
                            if info is not None:
                                for placeholder, image_path in rep["immagini"].items():
                                    if placeholder in info:
                                        saved = save_image(placeholder, image_path, "storage/immagini", width=0, height=0)
                                        with open(saved[info[1]], "rb") as img_file:
                                            image_bytes = img_file.read()

                                        run_img = new_par.add_run()
                                        run_img.add_picture(io.BytesIO(image_bytes), width=width, height=height)

                                        last_shape = doc.inline_shapes[-1]
                                        last_shape._inline.docPr.set('descr', new_alt or "")
                                        break
                            else:
                                # Copia immagine originale se nessun placeholder trovato
                                blip = shape._inline.graphic.graphicData.pic.blipFill.blip
                                image_part = doc.part.related_parts[blip.embed]
                                image_bytes = image_part.blob

                                run_img = new_par.add_run()
                                run_img.add_picture(io.BytesIO(image_bytes), width=width, height=height)

                                if alt_text:
                                    last_shape = doc.inline_shapes[-1]
                                    last_shape._inline.docPr.set('descr', alt_text)

                    nuovi_paragrafi.append(new_par)

            # Elimina blocco originale
            for idx in range(end_idx, start_idx - 1, -1):
                delete_paragraph(doc.paragraphs[idx])

            # Inserisci duplicati
            body = doc._body._element
            insert_pos = start_idx
            for new_par in nuovi_paragrafi:
                body.insert(insert_pos, new_par._element)
                insert_pos += 1

            i = insert_pos
            continue

        # -------------------------------
        # GESTIONE BLOCCO CONDIZIONALE {{if:ph:cond}}...{{fine_if:ph:cond}}
        # -------------------------------
        match_if = re.match(r"\{\{if:([^\}:]+):([^\}]*)\}\}", testo)
        if match_if:
            placeholder, condizione = match_if.group(1), match_if.group(2)
            ph = "{{" + placeholder + "}}"
            end_marker = f"{{{{fine_if:{placeholder}:{condizione}}}}}"
            end_idx = None
            for j in range(i+1, len(doc.paragraphs)):
                # if doc.paragraphs[j].text.endswith("{{fine_if:" + placeholder + ":" + condizione + "}}"):
                if end_marker in doc.paragraphs[j].text:
                    end_idx = j
                    break
            if end_idx is None:
                i += 1
                continue

            valore = None
            if ph in replacements:
                valore = replacements[ph]
            elif ph in image_replacements:
                valore = image_replacements[ph]

            soddisfatta = False
            if condizione == "*":
                soddisfatta = valore not in [None, ""]
            elif condizione == "":
                soddisfatta = valore == ""
            else:
                soddisfatta = valore == condizione

            doc.paragraphs[i].text = "<da_mantenere>" if soddisfatta else "<da_rimuovere>"
            doc.paragraphs[end_idx].text = "</da_mantenere>" if soddisfatta else "</da_rimuovere>"
            i = end_idx + 1
            continue

        i += 1

def process_file(file_path, replacements, image_replacements, replacements_for_each):
    # Verifica il tipo di file
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pptx":  
        all_replacements = copy.deepcopy(replacements)
        all_replacements.update(image_replacements)
        changed_presentation = os.path.join(UPLOAD_FOLDER, "changed.pptx")
        ppt = Presentation(file_path)
        # filtra_per(ppt, all_replacements)
        for_indexes = duplicate_and_replace_slide(ppt, replacements_for_each, replacements["{{numero_modelli}}"], replacements["{{numero_go}}"])
        replace_image_in_pptx(ppt, image_replacements)
        with open(changed_presentation, 'wb') as f:
            ppt.save(f)
        replacements_t = [(placeholder, text) for placeholder, text in replacements.items()]
        replacer = TextReplacer(changed_presentation, slides='', tables=True, charts=True, textframes=True)
        try:
            replacer.replace_text(replacements_t)
        except KeyError as e:
            print(f"Errore durante la sostituzione del testo: {e}", flush=True)
        replacer.write_presentation_to_file(changed_presentation)
        # for for_type, sequences in for_indexes.items():
        #     reps = replacements_for_each.get(for_type)
        #     for sequence in sequences:
        #         for ind, sliden in enumerate(sequence):
        #             slidenstr = str(sliden+1)
        #             replacer = TextReplacer(changed_presentation, slides=slidenstr, tables=True, charts=True, textframes=True)
        #             rep = reps[ind]['testuali']
        #             rep_t = [(placeholder, text) for placeholder, text in rep.items()]
        #             replacer.replace_text(rep_t)
        #             replacer.write_presentation_to_file(changed_presentation)
        for for_type, num_replace, sequences in for_indexes:
            reps = replacements_for_each.get(for_type)
            for ind, sliden in enumerate(sequences):
                slidenstr = str(sliden+1)
                replacer = TextReplacer(changed_presentation, slides=slidenstr, tables=True, charts=True, textframes=True)
                if reps is None or ind >= len(reps) or len(reps) == 0:
                    if for_type == "{{for_go:n}}":
                        reps = [{"testuali": {"{{gruppo_omogeneo_testuale_nome}}": "", "{{go_partecipanti}}": "", "{{go_adesione}}": "", "{{gruppo_omogeneo_rischio}}": ""}}]
                    elif for_type == "{{for_fg:n}}":
                        reps = [{"testuali": {"{{nome_fg}}": "", }}]
                    else:
                        continue
                for i in range(num_replace):
                    index = i + (ind * num_replace)
                    if index < len(reps):
                        rep = reps[index]['testuali']
                        rep_cambiato = {}
                        for key, value in rep.items():
                            k = "{{" + key.replace("{{", "").replace("}}", "") + ":" + str(i+1) + "}}"
                            rep_cambiato[k] = value
                        rep_t = [(placeholder, text) for placeholder, text in rep_cambiato.items()]
                        replacer.replace_text(rep_t)
                    else:
                        rep = reps[0]['testuali']
                        rep_cambiato = {}
                        for key, value in rep.items():
                            k = "{{" + key.replace("{{", "").replace("}}", "") + ":" + str(i+1) + "}}"
                            rep_cambiato[k] = ""
                        rep_t = [(placeholder, text) for placeholder, text in rep_cambiato.items()]
                        replacer.replace_text(rep_t)

                    replacer.write_presentation_to_file(changed_presentation)

        ppt = Presentation(changed_presentation)
        filtra_per(ppt, all_replacements)
        # Salva il file PPTX modificato
        ppt.save(changed_presentation)
    elif ext == ".docx":
        changed_presentation = os.path.join(UPLOAD_FOLDER, "changed.docx")
        docx_bytes = replace_text_in_docx(file_path, replacements, image_replacements, replacements_for_each)
        with open(changed_presentation, 'wb') as f:
            f.write(docx_bytes)
    else:   
        print("Tipo di file non supportato. Supportiamo solo .pptx e .docx.")

    with open(changed_presentation, 'rb') as f:
        file = f.read()
    return file
    
    # Verifica il tipo di file
    # ext = os.path.splitext(file_path)[1].lower()

    # for key, value in replacements.items():
    #     if key.startswith("{{if_tipo_woseq:") and key.endswith("}}"):
    #         if_tipo_woseq = key

    # if ext == ".pptx":
    #     ppt = Presentation(file_path)
    #     ppt = flitra_per_tipo_woseq(ppt, if_tipo_woseq)
    #     ppt = duplicate_and_replace_slide(ppt, replacements_for_each)
    #     ppt = replace_text_in_pptx(ppt, replacements, image_replacements)
    #     file_byte = salva_byte_pptx(ppt)
    #     return file_byte
    #     # replace_text_in_pptx(file_path, replacements, image_replacements)
    # elif ext == ".docx":
    #     file_byte = replace_text_in_docx(file_path, replacements, image_replacements)
    #     return file_byte
    # else:
    #     print("Tipo di file non supportato. Supportiamo solo .pptx e .docx.")

# # Esempio di utilizzo:
# replacements = {
#     "{{NOME_AZIENDA}}": "Mario Rossi",
#     "{{RUOLO}}": "Chef",
#     "{{azienda_prova}}": "MIND1234567890",
#     "{{TESTO_PROVA}}": "Testo di prova del titolo",
# 	"{{testo_prova_1}}": "Testo di prova del titolo",
# 	"{{testo_prova_2}}": "Testo di prova 2",
#     "Altro testo da cambiare": "Nuovo testo",
#     "XXX": "CIAO"
# }

# image_replacements = {
#      "{{immagine_prova1}}": "storage/immagini/image1.png",  # Inserisci il percorso dell'immagine da sostituire
# }

# # Percorso del file da modificare
# file_path = "storage/powerpoint1.pptx"  # O "path/to/your/file.docx"

# # Processare il file
# process_file(file_path, replacements, image_replacements)
